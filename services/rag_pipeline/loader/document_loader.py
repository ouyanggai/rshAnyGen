"""Document Loader - Support multiple file formats"""

import asyncio
from pathlib import Path
from typing import Dict, Any, List, Optional
import logging

logger = logging.getLogger(__name__)


class DocumentLoader:
    """Document loader supporting multiple formats"""

    TEXT_EXTENSIONS = {".txt", ".md", ".rst", ".log"}
    IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tiff"}
    WORD_EXTENSIONS = {".docx", ".doc"}
    EXCEL_EXTENSIONS = {".xlsx", ".xls"}
    POWERPOINT_EXTENSIONS = {".pptx", ".ppt"}
    PDF_EXTENSIONS = {".pdf"}

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """Initialize document loader

        Args:
            config: Configuration dictionary
        """
        self.config = config or {}

    async def load(self, file_path: str) -> Dict[str, Any]:
        """Load document content

        Args:
            file_path: Path to the file

        Returns:
            Dictionary containing document content and metadata

        Raises:
            ValueError: If file type is not supported
            FileNotFoundError: If file doesn't exist
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = path.suffix.lower()

        try:
            if ext in self.TEXT_EXTENSIONS:
                return await self._load_text(path, ext)
            elif ext in self.IMAGE_EXTENSIONS:
                return await self._load_image(path, ext)
            elif ext in self.WORD_EXTENSIONS:
                return await self._load_word(path, ext)
            elif ext in self.EXCEL_EXTENSIONS:
                return await self._load_excel(path, ext)
            elif ext in self.PDF_EXTENSIONS:
                return await self._load_pdf(path, ext)
            elif ext in self.POWERPOINT_EXTENSIONS:
                return await self._load_powerpoint(path, ext)
            elif ext in {".html", ".htm"}:
                return await self._load_html(path, ext)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
        except Exception as e:
            logger.error(f"Error loading file {file_path}: {e}")
            raise

    async def load_batch(self, file_paths: List[str]) -> List[Dict[str, Any]]:
        """Load multiple documents in batch

        Args:
            file_paths: List of file paths

        Returns:
            List of document dictionaries
        """
        tasks = [self.load(fp) for fp in file_paths]
        return await asyncio.gather(*tasks, return_exceptions=True)

    async def _load_text(self, path: Path, ext: str) -> Dict[str, Any]:
        """Load text-based files

        Args:
            path: Path to the file
            ext: File extension

        Returns:
            Document dictionary
        """
        try:
            content = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            # Try with alternative encodings
            for encoding in ["gbk", "gb2312", "latin-1"]:
                try:
                    content = path.read_text(encoding=encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise ValueError(f"Cannot decode file with any common encoding: {path}")

        return {
            "type": "text",
            "content": content,
            "metadata": {
                "format": ext,
                "file_name": path.name,
                "file_size": path.stat().st_size,
            },
        }

    async def _load_image(self, path: Path, ext: str) -> Dict[str, Any]:
        """Load image and extract text using OCR

        Args:
            path: Path to the image
            ext: File extension

        Returns:
            Document dictionary with OCR extracted text
        """
        try:
            from paddleocr import PaddleOCR

            # Initialize PaddleOCR (use angle classifier for rotated text)
            ocr = PaddleOCR(
                use_angle_cls=True,
                lang="ch",  # Chinese + English
                show_log=False,
            )

            # Run OCR
            result = ocr.ocr(str(path), cls=True)

            # Extract text from OCR result
            text_parts = []
            if result and result[0]:
                for line in result[0]:
                    if line and len(line) >= 2:
                        text_parts.append(line[1][0])

            content = "\n".join(text_parts)

            return {
                "type": "image",
                "content": content,
                "metadata": {
                    "format": ext,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "lines_detected": len(text_parts),
                },
            }
        except ImportError:
            logger.warning("paddleocr not installed. Run: pip install paddleocr")
            return {
                "type": "image",
                "content": "",
                "metadata": {
                    "format": ext,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "note": "paddleocr not installed",
                },
            }
        except Exception as e:
            logger.error(f"Error processing image with OCR: {e}")
            return {
                "type": "image",
                "content": "",
                "metadata": {
                    "format": ext,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "error": str(e),
                },
            }

    async def _load_word(self, path: Path, ext: str) -> Dict[str, Any]:
        """Load Word documents

        Args:
            path: Path to the Word document
            ext: File extension

        Returns:
            Document dictionary
        """
        try:
            from docx import Document

            doc = Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = "\n".join(paragraphs)

            return {
                "type": "word",
                "content": content,
                "metadata": {
                    "format": ext,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "paragraph_count": len(paragraphs),
                },
            }
        except ImportError:
            logger.warning("python-docx not installed. Run: pip install python-docx")
            return {
                "type": "word",
                "content": "",
                "metadata": {
                    "format": ext,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "note": "python-docx not installed",
                },
            }
        except Exception as e:
            logger.error(f"Error loading Word document: {e}")
            raise

    async def _load_excel(self, path: Path, ext: str) -> Dict[str, Any]:
        """Load Excel spreadsheets

        Args:
            path: Path to the Excel file
            ext: File extension

        Returns:
            Document dictionary
        """
        try:
            import openpyxl

            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            sheets_data = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    # Filter out completely empty rows
                    if any(cell is not None for cell in row):
                        rows.append([str(cell) if cell is not None else "" for cell in row])

                if rows:
                    # Convert to text representation
                    sheet_text = "\n".join(["\t".join(row) for row in rows])
                    sheets_data.append(f"=== Sheet: {sheet_name} ===\n{sheet_text}")

            content = "\n\n".join(sheets_data) if sheets_data else ""

            return {
                "type": "excel",
                "content": content,
                "metadata": {
                    "format": path.suffix,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "sheet_count": len(wb.sheetnames),
                },
            }
        except ImportError:
            logger.warning("openpyxl not installed. Run: pip install openpyxl")
            return {
                "type": "excel",
                "content": "",
                "metadata": {
                    "format": path.suffix,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "note": "openpyxl not installed",
                },
            }
        except Exception as e:
            logger.error(f"Error loading Excel file: {e}")
            raise

    async def _load_pdf(self, path: Path, ext: str) -> Dict[str, Any]:
        """Load PDF documents

        Args:
            path: Path to the PDF
            ext: File extension

        Returns:
            Document dictionary
        """
        try:
            import pypdf

            content_parts = []
            with open(path, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    text = page.extract_text()
                    if text.strip():
                        content_parts.append(text)

            content = "\n\n".join(content_parts)

            return {
                "type": "pdf",
                "content": content,
                "metadata": {
                    "format": ext,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "page_count": len(reader.pages),
                },
            }
        except ImportError:
            logger.warning("pypdf not installed. Run: pip install pypdf")
            return {
                "type": "pdf",
                "content": "",
                "metadata": {
                    "format": ext,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "note": "pypdf not installed",
                },
            }
        except Exception as e:
            logger.error(f"Error loading PDF: {e}")
            raise

    async def _load_powerpoint(self, path: Path, ext: str) -> Dict[str, Any]:
        """Load PowerPoint presentations

        Args:
            path: Path to the PowerPoint file
            ext: File extension

        Returns:
            Document dictionary
        """
        try:
            from pptx import Presentation

            prs = Presentation(path)
            slides_data = []

            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    slides_data.append(f"=== Slide {slide_idx + 1} ===\n" + "\n".join(slide_text))

            content = "\n\n".join(slides_data)

            return {
                "type": "powerpoint",
                "content": content,
                "metadata": {
                    "format": ext,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "slide_count": len(prs.slides),
                },
            }
        except ImportError:
            logger.warning("python-pptx not installed. Run: pip install python-pptx")
            return {
                "type": "powerpoint",
                "content": "",
                "metadata": {
                    "format": ext,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "note": "python-pptx not installed",
                },
            }
        except Exception as e:
            logger.error(f"Error loading PowerPoint file: {e}")
            raise

    async def _load_html(self, path: Path, ext: str) -> Dict[str, Any]:
        """Load HTML files

        Args:
            path: Path to the HTML file
            ext: File extension

        Returns:
            Document dictionary with extracted text
        """
        try:
            from bs4 import BeautifulSoup

            html_content = path.read_text(encoding="utf-8")
            soup = BeautifulSoup(html_content, "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            # Get text
            text = soup.get_text()

            # Break into lines and remove leading/trailing space
            lines = (line.strip() for line in text.splitlines())
            # Break multi-headlines into a line each
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            # Drop blank lines
            content = "\n".join(chunk for chunk in chunks if chunk)

            return {
                "type": "html",
                "content": content,
                "metadata": {
                    "format": ext,
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "title": soup.title.string if soup.title else "",
                },
            }
        except ImportError:
            logger.warning("beautifulsoup4 not installed. Run: pip install beautifulsoup4")
            # Fallback to basic text extraction
            try:
                import re

                html_content = path.read_text(encoding="utf-8")
                # Remove HTML tags
                content = re.sub(r"<[^>]+>", "", html_content)
                content = re.sub(r"\s+", " ", content).strip()

                return {
                    "type": "html",
                    "content": content,
                    "metadata": {
                        "format": ext,
                        "file_name": path.name,
                        "file_size": path.stat().st_size,
                        "note": "Using basic text extraction (beautifulsoup4 recommended)",
                    },
                }
            except Exception as e:
                logger.error(f"Error loading HTML file: {e}")
                raise
        except Exception as e:
            logger.error(f"Error loading HTML file: {e}")
            raise
